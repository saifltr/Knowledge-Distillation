#!/usr/bin/env python3
"""
Document Processor using Mistral OCR and AI

This script processes PDF, Excel, PowerPoint, and Word documents from a directory,
extracts text and images, and generates a comprehensive markdown file containing
all text and images with descriptions.
"""

import os
import io
import json
import base64
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime

# Mistral AI dependencies
from mistralai import Mistral
from mistralai import DocumentURLChunk, ImageURLChunk
from mistralai.models import OCRResponse

# Document processing dependencies
import fitz  # PyMuPDF for PDF processing
import openpyxl  # For Excel processing
from pptx import Presentation  # For PowerPoint processing
import docx  # For Word processing
from PIL import Image  # For image handling

def get_api_key() -> str:
    """Get Mistral API key from environment variable"""
    api_key = os.environ.get("MISTRAL_API_KEY", "6iAQWtFQnCfR7kKaIrbhS3ZmlCFscPbd")
    if not api_key:
        raise ValueError(
            "Mistral API key must be provided via MISTRAL_API_KEY environment variable"
        )
    return api_key

def process_pdf_with_mistral(client: Mistral, pdf_path: Path, model: str = "mistral-ocr-latest") -> OCRResponse:
    """Upload and process a PDF file with Mistral OCR"""
    print(f"Processing PDF with Mistral OCR: {pdf_path}")
    
    # Upload PDF file to Mistral's OCR service
    print("Uploading PDF file to Mistral OCR service...")
    uploaded_file = client.files.upload(
        file={
            "file_name": pdf_path.stem,
            "content": pdf_path.read_bytes(),
        },
        purpose="ocr",
    )
    
    # Get URL for the uploaded file
    signed_url = client.files.get_signed_url(file_id=uploaded_file.id, expiry=1)
    
    # Process PDF with OCR, including embedded images
    print("Processing PDF with OCR...")
    pdf_response = client.ocr.process(
        document=DocumentURLChunk(document_url=signed_url.url),
        model=model,
        include_image_base64=True
    )
    
    return pdf_response

def process_pdf_with_pymupdf(pdf_path: Path) -> Tuple[str, List[Dict[str, Any]]]:
    """Extract text and images from PDF using PyMuPDF"""
    print(f"Extracting content from PDF: {pdf_path}")
    
    document_text = []
    extracted_images = []
    
    # Open the PDF
    pdf_document = fitz.open(pdf_path)
    
    # Process each page
    for page_num, page in enumerate(pdf_document):
        # Extract text
        page_text = page.get_text()
        document_text.append(f"## Page {page_num + 1}\n\n{page_text}\n\n")
        
        # Extract images
        image_list = page.get_images(full=True)
        
        for img_index, img_info in enumerate(image_list):
            xref = img_info[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            
            # Convert image bytes to base64
            image_base64 = base64.b64encode(image_bytes).decode('utf-8')
            
            # Create unique image ID
            image_id = f"{pdf_path.stem}_page{page_num+1}_img{img_index+1}"
            
            extracted_images.append({
                "id": image_id,
                "base64": image_base64,
                "page": page_num + 1
            })
    
    pdf_document.close()
    return "\n".join(document_text), extracted_images

def extract_images_from_docx(docx_path: Path) -> List[Dict[str, Any]]:
    """Extract images from Word document"""
    print(f"Extracting images from Word document: {docx_path}")
    
    extracted_images = []
    doc = docx.Document(docx_path)
    
    image_index = 0
    # Extract images from document parts
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_bytes = rel.target_part.blob
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                
                # Create unique image ID
                image_id = f"{docx_path.stem}_img{image_index+1}"
                
                extracted_images.append({
                    "id": image_id,
                    "base64": image_base64,
                    "page": 1  # Word docs don't have pages in the same way
                })
                
                image_index += 1
            except Exception as e:
                print(f"Error extracting image: {e}")
    
    return extracted_images

def process_docx(docx_path: Path) -> Tuple[str, List[Dict[str, Any]]]:
    """Extract text and images from Word document"""
    print(f"Processing Word document: {docx_path}")
    
    doc = docx.Document(docx_path)
    document_text = []
    
    # Extract text paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            document_text.append(para.text)
    
    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text for cell in row.cells]
            document_text.append(" | ".join(row_text))
    
    # Get images
    extracted_images = extract_images_from_docx(docx_path)
    
    return "\n\n".join(document_text), extracted_images

def extract_images_from_pptx(pptx_path: Path) -> List[Dict[str, Any]]:
    """Extract images from PowerPoint presentation using python-pptx"""
    print(f"Extracting images from PowerPoint: {pptx_path}")
    
    try:
        # Import required modules
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        extracted_images = []
        
        # Open the presentation
        presentation = Presentation(pptx_path)
        
        # Function to iterate through all picture shapes in the presentation
        def iter_picture_shapes():
            for slide_num, slide in enumerate(presentation.slides, 1):
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        yield slide_num, shape_idx, shape
        
        # Process each picture shape
        for slide_num, shape_idx, picture in iter_picture_shapes():
            try:
                # Get image data
                image = picture.image
                image_bytes = image.blob
                
                # Convert to base64
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                
                # Create unique image ID
                image_id = f"{pptx_path.stem}_slide{slide_num}_img{shape_idx}"
                
                extracted_images.append({
                    "id": image_id,
                    "base64": image_base64,
                    "page": slide_num
                })
                
                print(f"Successfully extracted image {image_id} from slide {slide_num}")
            except Exception as e:
                print(f"Error extracting image from slide {slide_num}, shape {shape_idx}: {e}")
        
        print(f"Total images extracted from PowerPoint: {len(extracted_images)}")
        
        # If no images found, try the fallback method
        if not extracted_images:
            print("No images found using direct method, trying fallback...")
            return extract_images_from_pptx_fallback(pptx_path)
            
        return extracted_images
    
    except Exception as e:
        print(f"Error processing PowerPoint file: {e}")
        print("Trying fallback method...")
        return extract_images_from_pptx_fallback(pptx_path)

def extract_images_from_pptx_fallback(pptx_path: Path) -> List[Dict[str, Any]]:
    """Fallback method to extract images from PowerPoint using zipfile approach"""
    print(f"Using fallback method to extract images from PowerPoint: {pptx_path}")
    
    from zipfile import ZipFile
    import os
    
    extracted_images = []
    image_index = 0
    
    try:
        # PowerPoint files are essentially ZIP files with XML and media content
        with ZipFile(pptx_path) as pptx_zip:
            # Find all media files in the PowerPoint
            media_files = [item for item in pptx_zip.namelist() if item.startswith('ppt/media/')]
            
            if not media_files:
                print(f"No media files found in PowerPoint: {pptx_path}")
                return []
            
            print(f"Found {len(media_files)} media files in PowerPoint")
            
            # Process each media file
            for item in media_files:
                try:
                    # Extract the media file
                    image_bytes = pptx_zip.read(item)
                    
                    # Only process if it seems to be an image
                    if item.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.emf', '.wmf')):
                        image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                        
                        # Create unique image ID (using filename from PowerPoint)
                        filename = os.path.basename(item)
                        image_id = f"{pptx_path.stem}_img{image_index+1}_{filename}"
                        
                        extracted_images.append({
                            "id": image_id,
                            "base64": image_base64,
                            "page": 0  # We don't know which slide this belongs to
                        })
                        
                        image_index += 1
                        print(f"Extracted image: {image_id}")
                except Exception as e:
                    print(f"Error extracting media item {item}: {e}")
    
    except Exception as e:
        print(f"Error opening PowerPoint file: {e}")
    
    print(f"Total images extracted from PowerPoint (fallback method): {len(extracted_images)}")
    return extracted_images

def process_pptx(pptx_path: Path) -> Tuple[str, List[Dict[str, Any]]]:
    """Extract text and images from PowerPoint presentation"""
    print(f"Processing PowerPoint presentation: {pptx_path}")
    
    presentation = Presentation(pptx_path)
    document_text = []
    
    # Extract text from slides
    for slide_num, slide in enumerate(presentation.slides):
        slide_text = []
        slide_text.append(f"## Slide {slide_num + 1}")
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        
        document_text.append("\n\n".join(slide_text))
    
    # Get images
    extracted_images = extract_images_from_pptx(pptx_path)
    
    return "\n\n---\n\n".join(document_text), extracted_images

def process_excel(excel_path: Path) -> Tuple[str, List[Dict[str, Any]]]:
    """Extract text and potential images from Excel file"""
    print(f"Processing Excel file: {excel_path}")
    
    workbook = openpyxl.load_workbook(excel_path, data_only=True)
    document_text = []
    extracted_images = []
    
    # Process each worksheet
    for sheet_index, sheet in enumerate(workbook.worksheets):
        sheet_text = [f"## Sheet: {sheet.title}"]
        
        # Get the used range
        min_row, min_col = 1, 1
        max_row = sheet.max_row
        max_col = sheet.max_column
        
        # Create header row
        headers = []
        for col in range(min_col, max_col + 1):
            cell_value = sheet.cell(min_row, col).value
            headers.append(str(cell_value) if cell_value is not None else "")
        
        if any(header.strip() for header in headers):
            sheet_text.append("| " + " | ".join(headers) + " |")
            sheet_text.append("| " + " | ".join(["---"] * len(headers)) + " |")
        
        # Create data rows
        for row in range(min_row + 1, max_row + 1):
            row_values = []
            for col in range(min_col, max_col + 1):
                cell_value = sheet.cell(row, col).value
                row_values.append(str(cell_value) if cell_value is not None else "")
            
            # Only add rows with at least one non-empty cell
            if any(value.strip() for value in row_values):
                sheet_text.append("| " + " | ".join(row_values) + " |")
        
        # Extract images if present in Excel file
        for image in sheet._images:
            try:
                # Excel stores images in the drawing part
                image_bytes = image._data()
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                
                # Create unique image ID
                image_id = f"{excel_path.stem}_sheet{sheet_index+1}_img{len(extracted_images)+1}"
                
                extracted_images.append({
                    "id": image_id,
                    "base64": image_base64,
                    "page": sheet_index + 1
                })
            except Exception as e:
                print(f"Error extracting image from Excel: {e}")
        
        document_text.append("\n\n".join(sheet_text))
    
    return "\n\n---\n\n".join(document_text), extracted_images

def generate_image_description(client: Mistral, image_base64: str) -> str:
    """
    Generate a description for an image using Mistral AI's vision model
    
    Args:
        client: Mistral client
        image_base64: Base64-encoded image data
        
    Returns:
        Description of the image
    """
    try:
        # Create base64 data URL for the image - ensuring no double prefix
        if image_base64.startswith("data:image/jpeg;base64,"):
            base64_data_url = image_base64
        else:
            base64_data_url = f"data:image/jpeg;base64,{image_base64}"
        
        # Define the messages for the chat
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "Please describe this image in detail, focusing on what it shows and its visual characteristics:"
                    },
                    {
                        "type": "image_url",
                        "image_url": base64_data_url
                    }
                ]
            }
        ]
        
        # Get the chat response using the complete method
        chat_response = client.chat.complete(
            model="pixtral-12b-latest",  # Using pixtral for better visual analysis
            messages=messages
        )
        
        # Return the content of the response
        return chat_response.choices[0].message.content
        
    except Exception as e:
        print(f"Error generating image description: {e}")
        return "No description available"

def save_image_to_file(image_base64: str, image_id: str, output_dir: Path) -> str:
    """Save base64 image to a file and return the relative path"""
    try:
        # Create images directory if it doesn't exist
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Determine file extension (assuming JPEG for simplicity)
        file_extension = ".jpg"
        
        # Create a filename based on the image_id
        filename = f"{image_id}{file_extension}"
        image_path = output_dir / filename
        
        # Decode base64 and save to file
        image_data = base64.b64decode(image_base64)
        with open(image_path, 'wb') as f:
            f.write(image_data)
        
        return str(image_path)
    except Exception as e:
        print(f"Error saving image {image_id}: {e}")
        return None

def create_markdown_with_images(document_text: str, images: List[Dict[str, Any]], client: Mistral, images_dir: Path) -> str:
    """
    Create markdown content with images and descriptions
    
    Args:
        document_text: Text content extracted from document
        images: List of image information dictionaries
        client: Mistral client for generating image descriptions
        images_dir: Directory to save images
        
    Returns:
        Markdown string with document content and images with descriptions
    """
    markdown_content = document_text
    
    # Process each image and add to markdown
    for img_index, img_info in enumerate(images):
        img_id = img_info["id"]
        img_base64 = img_info["base64"]
        page_num = img_info.get("page", 1)
        
        # Generate image description
        print(f"Generating description for image {img_id}...")
        img_description = generate_image_description(client, img_base64)
        
        # Add delay to avoid rate limit issues
        import time
        time.sleep(1)  # Sleep for 1 second between API calls
        
        # Save image to file
        image_path = save_image_to_file(img_base64, img_id, images_dir)
        
        # Create relative path for markdown
        relative_path = os.path.relpath(image_path, images_dir.parent)
        
        # Create image reference
        image_markdown = f"![{img_id}]({relative_path})\n\n**Image Description:**\n\n{img_description}\n\n"
        
        # Find a place to insert the image in the markdown based on its page number
        page_marker = f"## Page {page_num}" if "Page" in document_text else f"## Slide {page_num}" if "Slide" in document_text else f"## Sheet"
        
        if page_marker in markdown_content:
            # Insert after the page/slide/sheet marker
            page_index = markdown_content.find(page_marker)
            next_marker_index = markdown_content.find("##", page_index + 1)
            
            if next_marker_index > -1:
                # Insert before the next marker
                insertion_point = next_marker_index
            else:
                # Insert at the end
                insertion_point = len(markdown_content)
            
            # Insert the image markdown
            markdown_content = markdown_content[:insertion_point] + "\n\n" + image_markdown + markdown_content[insertion_point:]
        else:
            # If no matching page marker, append at the end
            markdown_content += "\n\n" + image_markdown
    
    return markdown_content

def process_document(file_path: Path, client: Mistral) -> Tuple[str, List[Dict[str, Any]]]:
    """
    Process different document types and extract text and images
    
    Args:
        file_path: Path to the document
        client: Mistral client
        
    Returns:
        Tuple containing extracted text and list of images
    """
    file_extension = file_path.suffix.lower()
    
    try:
        if file_extension == ".pdf":
            # For PDFs, we can either use Mistral OCR or PyMuPDF
            # Here we'll use PyMuPDF for consistency with other file types
            return process_pdf_with_pymupdf(file_path)
            
        elif file_extension in [".xlsx", ".xls"]:
            return process_excel(file_path)
            
        elif file_extension in [".pptx", ".ppt"]:
            return process_pptx(file_path)
            
        elif file_extension in [".docx", ".doc"]:
            return process_docx(file_path)
            
        else:
            print(f"Unsupported file type: {file_extension}")
            return "", []
            
    except Exception as e:
        print(f"Error processing {file_path}: {e}")
        return "", []

def main():
    """Main function to process documents and generate markdown"""
    # Define input and output directories
    input_dir = Path("docs")
    output_dir = Path("markdown_output")
    
    # Create output directory if it doesn't exist
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # Create images directory
    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)
    
    # Get API key
    api_key = get_api_key()
    
    # Initialize Mistral client
    client = Mistral(api_key=api_key)
    
    # Check if input directory exists
    if not input_dir.exists():
        print(f"Input directory '{input_dir}' does not exist. Creating it...")
        input_dir.mkdir(parents=True, exist_ok=True)
        print(f"Please place your documents in the '{input_dir}' directory and run the script again.")
        return
    
    # Get list of documents to process
    supported_extensions = [".pdf", ".xlsx", ".xls", ".pptx", ".ppt", ".docx", ".doc"]
    document_files = [f for f in input_dir.iterdir() if f.is_file() and f.suffix.lower() in supported_extensions]
    
    if not document_files:
        print(f"No supported documents found in '{input_dir}' directory.")
        print(f"Supported formats: {', '.join(supported_extensions)}")
        return
    
    # Initialize markdown content
    current_time = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    markdown_sections = [
        "# Documents Processed with Mistral AI\n",
        f"*Processed on: {current_time}*\n\n",
        f"Total documents processed: {len(document_files)}\n\n",
        "---\n\n"
    ]
    
    # Process each document
    all_document_texts = []
    all_images = []
    
    for doc_path in document_files:
        print(f"\nProcessing document: {doc_path}")
        
        # Add document header to markdown
        doc_header = f"# Document: {doc_path.name}\n\n"
        all_document_texts.append(doc_header)
        
        # Process document
        document_text, images = process_document(doc_path, client)
        
        if document_text:
            all_document_texts.append(document_text)
        
        if images:
            # Add document info to each image
            for img in images:
                img["document"] = doc_path.name
            
            all_images.extend(images)
        
        # Add separator
        all_document_texts.append("\n\n---\n\n")
    
    # Combine all document texts
    combined_text = "\n".join(all_document_texts)
    
    # Create markdown with images and descriptions
    final_markdown = create_markdown_with_images(combined_text, all_images, client, images_dir)
    
    # Write to output file
    output_file = output_dir / "processed_documents.md"
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(final_markdown)
    
    print("\n\nProcessing completed!")
    print(f"Output markdown file: {output_file}")
    print(f"Total documents processed: {len(document_files)}")
    print(f"Total images processed: {len(all_images)}")
    print(f"Images saved to: {images_dir}")

if __name__ == "__main__":
    main()